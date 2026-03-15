"""
Extract text from PDF, DOCX, XLS, XLSX, CSV, PPTX and plain-text files.
Returns chunks with content_type for use by the indexer.
"""

import csv
import logging

logger = logging.getLogger(__name__)

# Chunking defaults (prose)
MAX_CONTENT_CHARS = 8000
CHUNK_OVERLAP = 200

# Prose extensions: read as plain text and chunk
PROSE_EXTENSIONS = {
    ".md", ".txt", ".py", ".js", ".ts", ".tsx", ".jsx", ".json",
    ".rst", ".yml", ".yaml", ".toml", ".cfg", ".ini", ".sh", ".bash", ".zsh",
}

# Extensions this module can extract (binary/structured)
EXTRACTOR_EXTENSIONS = {".pdf", ".docx", ".doc", ".xls", ".xlsx", ".csv", ".pptx"}


def _read_plain_text(path: str, max_size: int = 2_000_000) -> str | None:
    """Read file as UTF-8 text. Returns None on failure or binary."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read(max_size)
        if "\x00" in content:
            return None
        return content
    except Exception as e:
        logger.warning("Could not read %s as text: %s", path, e)
        return None


def _chunk_text(text: str, max_chars: int = MAX_CONTENT_CHARS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks for prose."""
    if not text or not text.strip():
        return []
    if len(text) <= max_chars:
        return [text.strip()]
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_chars
        chunk = text[start:end]
        if end < len(text):
            last_nl = chunk.rfind("\n")
            if last_nl > max_chars // 2:
                chunk = chunk[: last_nl + 1]
                end = start + last_nl + 1
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    return chunks


# --- PDF ---
def _extract_pdf(path: str) -> list[str] | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed; cannot extract PDF %s", path)
        return None
    try:
        reader = PdfReader(path)
        chunks = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                # One chunk per page (keeps tables coherent)
                chunks.append(text.strip())
        if not chunks:
            return None
        # If total content is small, merge pages to avoid tiny chunks
        merged = "\n\n".join(chunks)
        if len(merged) <= MAX_CONTENT_CHARS:
            return [merged]
        return _chunk_text(merged, MAX_CONTENT_CHARS, CHUNK_OVERLAP)
    except Exception as e:
        logger.warning("PDF extraction failed for %s: %s", path, e)
        return None


# --- DOCX ---
def _extract_docx(path: str) -> list[str] | None:
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed; cannot extract DOCX %s", path)
        return None
    try:
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        if not text.strip():
            return None
        return _chunk_text(text, MAX_CONTENT_CHARS, CHUNK_OVERLAP)
    except Exception as e:
        logger.warning("DOCX extraction failed for %s: %s", path, e)
        return None


# --- XLSX ---
def _extract_xlsx(path: str) -> list[str] | None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed; cannot extract XLSX %s", path)
        return None
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        chunks = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            header = [str(v) if v is not None else "" for v in rows[0]]
            header_line = " | ".join(header)
            chunk_lines = [f"Sheet: {sheet_name}\nHeader: {header_line}"]
            for row in rows[1:]:
                values = [str(v) if v is not None else "" for v in row]
                chunk_lines.append(" | ".join(values))
            text = "\n".join(chunk_lines)
            if len(text) <= MAX_CONTENT_CHARS:
                chunks.append(text)
            else:
                # Chunk by row groups, keeping header in each
                current = [f"Sheet: {sheet_name}\nHeader: {header_line}"]
                current_len = len(current[0])
                for row in rows[1:]:
                    values = [str(v) if v is not None else "" for v in row]
                    line = " | ".join(values)
                    if current_len + len(line) + 1 > MAX_CONTENT_CHARS and len(current) > 1:
                        chunks.append("\n".join(current))
                        current = [f"Sheet: {sheet_name}\nHeader: {header_line}", line]
                        current_len = len(current[0]) + len(line)
                    else:
                        current.append(line)
                        current_len += len(line) + 1
                if len(current) > 1:
                    chunks.append("\n".join(current))
        wb.close()
        return chunks if chunks else None
    except Exception as e:
        logger.warning("XLSX extraction failed for %s: %s", path, e)
        return None


# --- XLS (legacy) ---
def _extract_xls(path: str) -> list[str] | None:
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd not installed; cannot extract XLS %s", path)
        return None
    try:
        wb = xlrd.open_workbook(path)
        chunks = []
        for sheet in wb.sheets():
            if sheet.nrows == 0:
                continue
            header = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
            header_line = " | ".join(header)
            chunk_lines = [f"Sheet: {sheet.name}\nHeader: {header_line}"]
            for r in range(1, sheet.nrows):
                values = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                chunk_lines.append(" | ".join(values))
            text = "\n".join(chunk_lines)
            if len(text) <= MAX_CONTENT_CHARS:
                chunks.append(text)
            else:
                current = [f"Sheet: {sheet.name}\nHeader: {header_line}"]
                current_len = len(current[0])
                for r in range(1, sheet.nrows):
                    values = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                    line = " | ".join(values)
                    if current_len + len(line) + 1 > MAX_CONTENT_CHARS and len(current) > 1:
                        chunks.append("\n".join(current))
                        current = [f"Sheet: {sheet.name}\nHeader: {header_line}", line]
                        current_len = len(current[0]) + len(line)
                    else:
                        current.append(line)
                        current_len += len(line) + 1
                if len(current) > 1:
                    chunks.append("\n".join(current))
        return chunks if chunks else None
    except Exception as e:
        logger.warning("XLS extraction failed for %s: %s", path, e)
        return None


# --- CSV ---
def _extract_csv(path: str) -> list[str] | None:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = list(reader)
        if not rows:
            return None
        header = rows[0]
        header_line = " | ".join(header)
        chunk_lines = [f"Header: {header_line}"]
        for row in rows[1:]:
            chunk_lines.append(" | ".join(str(c) for c in row))
        text = "\n".join(chunk_lines)
        if len(text) <= MAX_CONTENT_CHARS:
            return [text]
        chunks = []
        current = [f"Header: {header_line}"]
        current_len = len(current[0])
        for row in rows[1:]:
            line = " | ".join(str(c) for c in row)
            if current_len + len(line) + 1 > MAX_CONTENT_CHARS and len(current) > 1:
                chunks.append("\n".join(current))
                current = [f"Header: {header_line}", line]
                current_len = len(current[0]) + len(line)
            else:
                current.append(line)
                current_len += len(line) + 1
        if len(current) > 1:
            chunks.append("\n".join(current))
        return chunks
    except Exception as e:
        logger.warning("CSV extraction failed for %s: %s", path, e)
        return None


# --- PPTX ---
def _extract_pptx(path: str) -> list[str] | None:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; cannot extract PPTX %s", path)
        return None
    try:
        prs = Presentation(path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            parts = [f"Slide {i + 1}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())
            text = "\n".join(parts)
            if text.strip() and text != f"Slide {i + 1}:":
                chunks.append(text)
        return chunks if chunks else None
    except Exception as e:
        logger.warning("PPTX extraction failed for %s: %s", path, e)
        return None


def extract_content(file_path: str, ext: str) -> tuple[list[str], str] | None:
    """
    Extract text from a file and return (list of chunks, content_type).
    content_type is one of: "prose", "table", "slide".
    Returns None if the file cannot be read or the format is unsupported.
    """
    ext_lower = ext.lower() if ext.startswith(".") else f".{ext.lower()}"

    if ext_lower in PROSE_EXTENSIONS:
        content = _read_plain_text(file_path)
        if not content or not content.strip():
            return None
        chunks = _chunk_text(content, MAX_CONTENT_CHARS, CHUNK_OVERLAP)
        return (chunks, "prose") if chunks else None

    if ext_lower == ".pdf":
        chunks = _extract_pdf(file_path)
        return (chunks, "prose") if chunks else None
    if ext_lower == ".docx":
        chunks = _extract_docx(file_path)
        return (chunks, "prose") if chunks else None
    if ext_lower == ".doc":
        # Legacy .doc: python-docx does not support it. Try plain text (works only if saved as text).
        content = _read_plain_text(file_path)
        chunks = _chunk_text(content, MAX_CONTENT_CHARS, CHUNK_OVERLAP) if content else None
        return (chunks, "prose") if chunks else None
    if ext_lower == ".xlsx":
        chunks = _extract_xlsx(file_path)
        return (chunks, "table") if chunks else None
    if ext_lower == ".xls":
        chunks = _extract_xls(file_path)
        return (chunks, "table") if chunks else None
    if ext_lower == ".csv":
        chunks = _extract_csv(file_path)
        return (chunks, "table") if chunks else None
    if ext_lower == ".pptx":
        chunks = _extract_pptx(file_path)
        return (chunks, "slide") if chunks else None

    return None
