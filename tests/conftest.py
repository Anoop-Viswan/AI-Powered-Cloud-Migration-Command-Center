"""Pytest fixtures: temp project dir with sample files for extractors and semantic_search."""

# Disable admin auth during tests so admin/diagnostics endpoints are reachable without login
import os
os.environ.pop("ADMIN_USERNAME", None)
os.environ.pop("ADMIN_PASSWORD", None)

import json
from pathlib import Path

import pytest


def _write_minimal_pdf(path: Path) -> None:
    """Write a minimal valid PDF (blank page). Extractors may return None for no text."""
    from pypdf import PdfWriter
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.write(path)


def _write_minimal_docx(path: Path) -> None:
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from DOCX unit test.")
    doc.save(str(path))


def _write_minimal_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 42])
    ws.append(["Beta", 100])
    wb.save(path)


def _write_minimal_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tb.text_frame.text = "Hello from PPTX unit test."
    prs.save(str(path))


@pytest.fixture
def temp_project_dir(tmp_path):
    """A temporary directory with sample files and optional manifest."""
    root = Path(tmp_path)
    # Root-level file -> application "default"
    (root / "readme.txt").write_text("Project readme at root.", encoding="utf-8")
    # App folder
    app_dir = root / "MyApp"
    app_dir.mkdir()
    (app_dir / "doc.txt").write_text("App document content.", encoding="utf-8")
    (app_dir / "data.csv").write_text("col1,col2\n1,2\na,b", encoding="utf-8")
    yield root


@pytest.fixture
def temp_project_with_manifest(temp_project_dir):
    """Same as temp_project_dir plus manifest.json."""
    manifest = {
        "MyApp": {"technology": "Python", "tools": "Docker", "description": "Test app"},
        "default": {"description": "Root files"},
    }
    (temp_project_dir / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    yield temp_project_dir


@pytest.fixture
def sample_txt_path(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text("Plain text content for unit test.\nSecond line.", encoding="utf-8")
    return path


@pytest.fixture
def sample_csv_path(tmp_path):
    path = tmp_path / "sample.csv"
    path.write_text("Project,Owner\nAlpha,Alice\nBeta,Bob", encoding="utf-8")
    return path


@pytest.fixture
def sample_pdf_path(tmp_path):
    path = tmp_path / "sample.pdf"
    _write_minimal_pdf(path)
    return path


@pytest.fixture
def sample_docx_path(tmp_path):
    path = tmp_path / "sample.docx"
    _write_minimal_docx(path)
    return path


@pytest.fixture
def sample_xlsx_path(tmp_path):
    path = tmp_path / "sample.xlsx"
    _write_minimal_xlsx(path)
    return path


@pytest.fixture
def sample_pptx_path(tmp_path):
    path = tmp_path / "sample.pptx"
    _write_minimal_pptx(path)
    return path
